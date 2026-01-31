"""
MindSparkle Document Intelligence Service v2.0
Python-based extraction with production-grade libraries

Architecture:
- PDF: Google Document AI → Adobe PDF (if confidence < 85%) → OCR fallback
- PPTX: python-pptx + speaker notes + tables + OCR embedded images
- DOCX: python-docx + OCR embedded images  
- TXT: Direct read
"""

import os
import io
import uuid
import time
import json
import base64
import tempfile
import requests
from typing import Optional, Dict, Any, List
from flask import Flask, request, jsonify
from flask_cors import CORS

import google.auth
from google.auth.transport.requests import Request as GoogleAuthRequest

# Document processing libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
import fitz  # PyMuPDF for PDF
from PIL import Image

# Google Cloud
from google.cloud import documentai_v1 as documentai
from google.cloud import vision

app = Flask(__name__)
CORS(app)

# Configuration
GOOGLE_PROJECT_ID = (
    os.environ.get('GOOGLE_PROJECT_ID')
    or os.environ.get('GOOGLE_CLOUD_PROJECT')
    or os.environ.get('GOOGLE_PROJECT')
    or 'mindsparkle'
)
DOCUMENT_AI_LOCATION = os.environ.get('DOCUMENT_AI_LOCATION', 'us')
DOCUMENT_AI_PROCESSOR_ID = os.environ.get('DOCUMENT_AI_PROCESSOR_ID', '')
OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY', '')

VERTEX_LOCATION = os.environ.get('VERTEX_LOCATION', os.environ.get('GOOGLE_CLOUD_REGION', 'us-central1'))
VERTEX_GEMINI_OCR_ENABLED = os.environ.get('VERTEX_GEMINI_OCR_ENABLED', 'false').lower() in ['1', 'true', 'yes']
VERTEX_GEMINI_MODEL = os.environ.get('VERTEX_GEMINI_MODEL', 'gemini-1.5-pro')

# =============================================================================
# HEALTH CHECK
# =============================================================================

@app.route('/', methods=['GET'])
def index():
    return jsonify({
        'service': 'MindSparkle Document Intelligence',
        'version': '2.0.0',
        'status': 'healthy',
        'architecture': 'Python-based extraction → Canonical Model → AI Understanding',
        'endpoints': {
            'POST /extract': 'Extract document to canonical model',
            'POST /understand': 'Get AI understanding of document',
            'POST /process': 'Full pipeline: extract + understand',
            'GET /health': 'Health check'
        },
        'extractors': {
            'pdf': 'Google Document AI → Adobe fallback → OCR',
            'pptx': 'python-pptx with speaker notes, tables, image OCR',
            'docx': 'python-docx with image OCR',
            'txt': 'Direct read'
        }
    })

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'timestamp': time.strftime('%Y-%m-%dT%H:%M:%SZ')})

# =============================================================================
# EXTRACT ENDPOINT
# =============================================================================

@app.route('/extract', methods=['POST'])
def extract():
    start_time = time.time()
    request_id = str(uuid.uuid4())[:8]
    
    try:
        # Get file from request
        file_buffer = None
        filename = 'document'
        mimetype = ''
        
        # Check for file upload
        if 'file' in request.files:
            file = request.files['file']
            file_buffer = file.read()
            filename = file.filename or 'document'
            mimetype = file.content_type or ''
            print(f"[{request_id}] Received file upload: {filename} ({len(file_buffer)} bytes)")
        
        # Check for signed URL in JSON body
        elif request.is_json and request.json.get('signedUrl'):
            data = request.json
            signed_url = data['signedUrl']
            filename = data.get('fileName', 'document')
            mimetype = data.get('mimeType', '')
            
            print(f"[{request_id}] Downloading from signed URL: {filename}")
            response = requests.get(signed_url, timeout=300)
            response.raise_for_status()
            file_buffer = response.content
            print(f"[{request_id}] Downloaded {len(file_buffer)} bytes")
        
        else:
            return jsonify({'error': 'No file provided'}), 400
        
        # Detect file type
        file_type = detect_file_type(filename, mimetype)
        print(f"[{request_id}] Extracting {file_type}: {filename}")
        
        # Extract based on type
        if file_type == 'pdf':
            result = extract_pdf(file_buffer, request_id)
        elif file_type in ['pptx', 'ppt']:
            result = extract_pptx(file_buffer, request_id)
        elif file_type in ['docx', 'doc']:
            result = extract_docx(file_buffer, request_id)
        elif file_type == 'txt':
            result = extract_txt(file_buffer, request_id)
        else:
            return jsonify({'error': f'Unsupported file type: {file_type}'}), 400
        
        # Build canonical model
        canonical = build_canonical_model(result, filename, file_type, len(file_buffer), request_id)
        
        processing_time = int((time.time() - start_time) * 1000)
        print(f"[{request_id}] Extraction complete in {processing_time}ms")
        
        return jsonify({
            'success': True,
            'requestId': request_id,
            'processingTime': processing_time,
            'canonical': canonical,
            'metadata': {
                'filename': filename,
                'fileType': file_type,
                'fileSize': len(file_buffer),
                'extractionMethod': result.get('method', 'unknown'),
                'confidence': result.get('confidence', 0)
            }
        })
        
    except Exception as e:
        print(f"[{request_id}] Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e), 'requestId': request_id}), 500

# =============================================================================
# PDF EXTRACTION - Google Document AI → Adobe fallback → OCR
# =============================================================================

MAX_PAGE_THUMBNAILS = int(os.environ.get('MAX_PAGE_THUMBNAILS', '25'))
THUMBNAIL_MAX_WIDTH = int(os.environ.get('THUMBNAIL_MAX_WIDTH', '900'))
THUMBNAIL_JPEG_QUALITY = int(os.environ.get('THUMBNAIL_JPEG_QUALITY', '55'))


def _render_page_thumbnail_data_url(page, request_id: str) -> Optional[str]:
    """Render a lightweight JPEG thumbnail for a PDF page and return a data URL."""
    try:
        # Compute a scale so the rendered pixmap is roughly THUMBNAIL_MAX_WIDTH.
        page_width = float(page.rect.width) if page and page.rect else 0.0
        if page_width <= 0:
            scale = 1.0
        else:
            scale = min(2.0, max(0.25, THUMBNAIL_MAX_WIDTH / page_width))

        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        mode = "RGB" if pix.n >= 3 else "L"
        img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
        if mode != "RGB":
            img = img.convert("RGB")

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=THUMBNAIL_JPEG_QUALITY, optimize=True)
        b64 = base64.b64encode(buf.getvalue()).decode('utf-8')
        return f"data:image/jpeg;base64,{b64}"
    except Exception as e:
        print(f"[{request_id}] Thumbnail render failed: {e}")
        return None

def extract_pdf(file_buffer: bytes, request_id: str) -> Dict[str, Any]:
    """Extract text from PDF using Document AI with fallbacks"""
    
    text_blocks = []
    tables = []
    images = []
    raw_text = ''
    method = 'unknown'
    confidence = 0.0
    page_count = 1
    
    # Try Google Document AI first
    if DOCUMENT_AI_PROCESSOR_ID:
        try:
            print(f"[{request_id}] Trying Google Document AI...")
            result = extract_pdf_with_document_ai(file_buffer, request_id)
            if result and result.get('confidence', 0) >= 0.7:
                print(f"[{request_id}] Document AI success: {result.get('confidence', 0)*100:.0f}% confidence")
                return result
            else:
                print(f"[{request_id}] Document AI low confidence, trying fallback...")
        except Exception as e:
            print(f"[{request_id}] Document AI failed: {e}")
    
    # Fallback to PyMuPDF (local extraction)
    try:
        print(f"[{request_id}] Using PyMuPDF extraction...")
        doc = fitz.open(stream=file_buffer, filetype="pdf")
        page_count = doc.page_count or 1
        thumbnails_added = 0
        
        for page_num, page in enumerate(doc):
            page_text = page.get_text("text")
            if page_text.strip():
                text_blocks.append({
                    'text': page_text.strip(),
                    'page': page_num + 1,
                    'confidence': 0.85,
                    'type': 'paragraph'
                })
            
            # Extract images for OCR if text is sparse
            if len(page_text.strip()) < 100:
                for img_index, img in enumerate(page.get_images()):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha > 3:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix.tobytes("png")
                        
                        # OCR the image
                        ocr_text = ocr_image(img_data, request_id)
                        if ocr_text:
                            text_blocks.append({
                                'text': ocr_text,
                                'page': page_num + 1,
                                'confidence': 0.7,
                                'type': 'ocr'
                            })
                    except:
                        pass

            # If we still have no meaningful text for this page, attach a thumbnail
            if thumbnails_added < MAX_PAGE_THUMBNAILS:
                try:
                    # Determine if any block exists for this page
                    has_page_text = any(
                        (b.get('page') == page_num + 1 and len(str(b.get('text', '')).strip()) >= 20)
                        for b in text_blocks
                    )
                    if not has_page_text:
                        data_url = _render_page_thumbnail_data_url(page, request_id)
                        if data_url:
                            images.append({
                                'page': page_num + 1,
                                'type': 'page_thumbnail',
                                'data_url': data_url,
                                'caption': 'Page thumbnail (no text extracted)'
                            })
                            thumbnails_added += 1
                except Exception:
                    pass
        
        doc.close()
        
        raw_text = '\n\n'.join([b['text'] for b in text_blocks])
        confidence = 0.85 if raw_text else 0.0
        method = 'pymupdf'
        
        # If still no text, try full page OCR
        if len(raw_text) < 100:
            print(f"[{request_id}] Low text extraction, trying full page OCR...")
            doc = fitz.open(stream=file_buffer, filetype="pdf")
            page_count = doc.page_count or page_count
            for page_num, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                ocr_text = ocr_image(img_data, request_id)
                if ocr_text:
                    text_blocks.append({
                        'text': ocr_text,
                        'page': page_num + 1,
                        'confidence': 0.6,
                        'type': 'page_ocr'
                    })
            doc.close()
            raw_text = '\n\n'.join([b['text'] for b in text_blocks])
            method = 'pymupdf_ocr'
            confidence = 0.6
        
    except Exception as e:
        print(f"[{request_id}] PyMuPDF failed: {e}")
    
    return {
        'method': method,
        'confidence': confidence,
        'text_blocks': text_blocks,
        'tables': tables,
        'images': images,
        'page_count': page_count,
        'raw_text': raw_text
    }


def extract_pdf_with_document_ai(file_buffer: bytes, request_id: str) -> Dict[str, Any]:
    """Use Google Document AI for high-quality PDF extraction"""
    
    client = documentai.DocumentProcessorServiceClient()
    name = f"projects/{GOOGLE_PROJECT_ID}/locations/{DOCUMENT_AI_LOCATION}/processors/{DOCUMENT_AI_PROCESSOR_ID}"
    
    raw_document = documentai.RawDocument(content=file_buffer, mime_type="application/pdf")
    request = documentai.ProcessRequest(name=name, raw_document=raw_document)
    
    result = client.process_document(request=request)
    document = result.document
    
    text_blocks = []
    tables = []
    
    # Extract text blocks
    for page in document.pages:
        page_num = page.page_number
        
        for paragraph in page.paragraphs:
            text = get_text_from_layout(paragraph.layout, document.text)
            confidence = paragraph.layout.confidence if paragraph.layout.confidence else 0.9
            text_blocks.append({
                'text': text,
                'page': page_num,
                'confidence': confidence,
                'type': 'paragraph'
            })
        
        # Extract tables
        for table in page.tables:
            table_data = []
            for row in table.header_rows + table.body_rows:
                row_data = []
                for cell in row.cells:
                    cell_text = get_text_from_layout(cell.layout, document.text)
                    row_data.append(cell_text)
                table_data.append(row_data)
            tables.append({
                'rows': table_data,
                'page': page_num,
                'confidence': 0.9
            })
    
    raw_text = document.text
    avg_confidence = sum(b['confidence'] for b in text_blocks) / len(text_blocks) if text_blocks else 0.9
    
    return {
        'method': 'document_ai',
        'confidence': avg_confidence,
        'text_blocks': text_blocks,
        'tables': tables,
        'images': [],
        'page_count': len(document.pages),
        'raw_text': raw_text
    }


def get_text_from_layout(layout, full_text: str) -> str:
    """Extract text from Document AI layout"""
    text = ""
    for segment in layout.text_anchor.text_segments:
        start = int(segment.start_index) if segment.start_index else 0
        end = int(segment.end_index)
        text += full_text[start:end]
    return text.strip()

# =============================================================================
# PPTX EXTRACTION - python-pptx with notes, tables, image OCR
# =============================================================================

def extract_pptx(file_buffer: bytes, request_id: str) -> Dict[str, Any]:
    """Extract text from PPTX using python-pptx"""
    
    print(f"[{request_id}] Extracting PPTX with python-pptx...")
    
    text_blocks = []
    tables = []
    images = []
    
    try:
        prs = Presentation(io.BytesIO(file_buffer))
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                            text_blocks.append({
                                'text': text,
                                'page': slide_num,
                                'confidence': 0.95,
                                'type': 'text'
                            })
                
                # Tables
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip() if cell.text else ''
                            row_data.append(cell_text)
                        table_data.append(row_data)
                    tables.append({
                        'rows': table_data,
                        'page': slide_num,
                        'confidence': 0.9
                    })
                    # Also add table text to blocks
                    table_text = '\n'.join([' | '.join(row) for row in table_data])
                    if table_text.strip():
                        text_blocks.append({
                            'text': f"[Table]\n{table_text}",
                            'page': slide_num,
                            'confidence': 0.9,
                            'type': 'table'
                        })
                
                # Images - OCR if slide has little text
                if hasattr(shape, 'image') and len(' '.join(slide_texts)) < 50:
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = ocr_image(image_bytes, request_id)
                        if ocr_text and len(ocr_text) > 20:
                            text_blocks.append({
                                'text': ocr_text,
                                'page': slide_num,
                                'confidence': 0.7,
                                'type': 'image_ocr'
                            })
                            images.append({
                                'page': slide_num,
                                'ocr_text': ocr_text
                            })
                    except Exception as e:
                        print(f"[{request_id}] Image OCR failed: {e}")
            
            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_blocks.append({
                        'text': f"[Speaker Notes]\n{notes_text}",
                        'page': slide_num,
                        'confidence': 0.95,
                        'type': 'speaker_notes'
                    })
        
        raw_text = '\n\n'.join([b['text'] for b in text_blocks])
        
        print(f"[{request_id}] Extracted {len(prs.slides)} slides, {len(text_blocks)} blocks, {len(tables)} tables")
        
        return {
            'method': 'python_pptx',
            'confidence': 0.95,
            'text_blocks': text_blocks,
            'tables': tables,
            'images': images,
            'page_count': len(prs.slides),
            'raw_text': raw_text
        }
        
    except Exception as e:
        print(f"[{request_id}] python-pptx failed: {e}")
        raise

# =============================================================================
# DOCX EXTRACTION - python-docx with image OCR
# =============================================================================

def extract_docx(file_buffer: bytes, request_id: str) -> Dict[str, Any]:
    """Extract text from DOCX using python-docx"""
    
    print(f"[{request_id}] Extracting DOCX with python-docx...")
    
    text_blocks = []
    tables = []
    images = []
    
    try:
        doc = Document(io.BytesIO(file_buffer))
        
        # Extract paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Detect headers
                para_type = 'heading' if para.style.name.startswith('Heading') else 'paragraph'
                text_blocks.append({
                    'text': text,
                    'page': 1,  # DOCX doesn't have page info
                    'confidence': 0.95,
                    'type': para_type
                })
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            tables.append({
                'rows': table_data,
                'page': 1,
                'confidence': 0.9
            })
            # Also add table text to blocks
            table_text = '\n'.join([' | '.join(row) for row in table_data])
            if table_text.strip():
                text_blocks.append({
                    'text': f"[Table]\n{table_text}",
                    'page': 1,
                    'confidence': 0.9,
                    'type': 'table'
                })
        
        # Extract images and OCR them
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    ocr_text = ocr_image(image_data, request_id)
                    if ocr_text and len(ocr_text) > 20:
                        text_blocks.append({
                            'text': f"[Image Text]\n{ocr_text}",
                            'page': 1,
                            'confidence': 0.7,
                            'type': 'image_ocr'
                        })
                        images.append({'ocr_text': ocr_text})
                except Exception as e:
                    print(f"[{request_id}] Image OCR failed: {e}")
        
        raw_text = '\n\n'.join([b['text'] for b in text_blocks])
        
        print(f"[{request_id}] Extracted {len(text_blocks)} blocks, {len(tables)} tables")
        
        return {
            'method': 'python_docx',
            'confidence': 0.95,
            'text_blocks': text_blocks,
            'tables': tables,
            'images': images,
            'page_count': 1,
            'raw_text': raw_text
        }
        
    except Exception as e:
        print(f"[{request_id}] python-docx failed: {e}")
        raise

# =============================================================================
# TXT EXTRACTION
# =============================================================================

def extract_txt(file_buffer: bytes, request_id: str) -> Dict[str, Any]:
    """Extract text from plain text file"""
    
    # Try different encodings
    text = None
    for encoding in ['utf-8', 'latin-1', 'cp1252']:
        try:
            text = file_buffer.decode(encoding)
            break
        except:
            continue
    
    if not text:
        text = file_buffer.decode('utf-8', errors='ignore')
    
    return {
        'method': 'direct_read',
        'confidence': 1.0,
        'text_blocks': [{'text': text, 'page': 1, 'confidence': 1.0, 'type': 'text'}],
        'tables': [],
        'images': [],
        'page_count': 1,
        'raw_text': text
    }

# =============================================================================
# OCR HELPER - Google Cloud Vision
# =============================================================================

def ocr_image(image_data: bytes, request_id: str) -> Optional[str]:
    """OCR an image using Google Cloud Vision"""
    
    try:
        client = vision.ImageAnnotatorClient()
        image = vision.Image(content=image_data)
        response = client.text_detection(image=image)
        
        if response.text_annotations:
            return response.text_annotations[0].description.strip()
        # Vision returned no text. Try stronger multimodal fallback(s).
        if VERTEX_GEMINI_OCR_ENABLED:
            try:
                text = ocr_with_vertex_gemini(image_data, request_id)
                if text:
                    return text
            except Exception as e:
                print(f"[{request_id}] Vertex Gemini OCR failed: {e}")

        if OPENAI_API_KEY:
            try:
                return ocr_with_openai(image_data, request_id)
            except Exception as e:
                print(f"[{request_id}] OpenAI OCR fallback failed: {e}")

        return None
        
    except Exception as e:
        print(f"[{request_id}] Vision OCR failed: {e}")

        # Fallback to Vertex Gemini if available (uses service account)
        if VERTEX_GEMINI_OCR_ENABLED:
            try:
                text = ocr_with_vertex_gemini(image_data, request_id)
                if text:
                    return text
            except Exception as e:
                print(f"[{request_id}] Vertex Gemini OCR failed: {e}")
        
        # Fallback to OpenAI Vision if available
        if OPENAI_API_KEY:
            try:
                return ocr_with_openai(image_data, request_id)
            except:
                pass
        
        return None


def ocr_with_vertex_gemini(image_data: bytes, request_id: str) -> Optional[str]:
    """OCR using Vertex AI Gemini (service account auth) as fallback."""

    # Acquire OAuth token for Vertex AI
    creds, _ = google.auth.default(scopes=['https://www.googleapis.com/auth/cloud-platform'])
    creds.refresh(GoogleAuthRequest())
    token = creds.token
    if not token:
        raise RuntimeError('Failed to acquire Google auth token')

    base64_image = base64.b64encode(image_data).decode('utf-8')
    endpoint = (
        f"https://{VERTEX_LOCATION}-aiplatform.googleapis.com/v1/projects/{GOOGLE_PROJECT_ID}"
        f"/locations/{VERTEX_LOCATION}/publishers/google/models/{VERTEX_GEMINI_MODEL}:generateContent"
    )

    payload = {
        'contents': [
            {
                'role': 'user',
                'parts': [
                    {'text': 'Extract ALL text from this image. Output only the text, nothing else.'},
                    {'inlineData': {'mimeType': 'image/png', 'data': base64_image}},
                ],
            }
        ],
        'generationConfig': {
            'temperature': 0.0,
            'maxOutputTokens': 8192,
        },
    }

    resp = requests.post(
        endpoint,
        headers={
            'Authorization': f'Bearer {token}',
            'Content-Type': 'application/json',
        },
        json=payload,
        timeout=30,
    )

    if resp.status_code != 200:
        raise RuntimeError(f"Vertex Gemini error {resp.status_code}: {resp.text[:500]}")

    data = resp.json()
    parts = (((data.get('candidates') or [{}])[0].get('content') or {}).get('parts') or [])
    text = ''.join([p.get('text', '') for p in parts if isinstance(p, dict)]).strip()
    return text or None


def ocr_with_openai(image_data: bytes, request_id: str) -> Optional[str]:
    """OCR using OpenAI Vision API as fallback"""
    
    base64_image = base64.b64encode(image_data).decode('utf-8')
    
    response = requests.post(
        'https://api.openai.com/v1/chat/completions',
        headers={
            'Authorization': f'Bearer {OPENAI_API_KEY}',
            'Content-Type': 'application/json'
        },
        json={
            'model': 'gpt-4o',
            'messages': [{
                'role': 'user',
                'content': [
                    {'type': 'text', 'text': 'Extract ALL text from this image. Output only the text, nothing else.'},
                    {'type': 'image_url', 'image_url': {'url': f'data:image/png;base64,{base64_image}'}}
                ]
            }],
            'max_tokens': 4000
        },
        timeout=30
    )
    
    if response.status_code == 200:
        return response.json()['choices'][0]['message']['content'].strip()
    return None

# =============================================================================
# CANONICAL MODEL
# =============================================================================

def build_canonical_model(result: Dict, filename: str, file_type: str, file_size: int, request_id: str) -> Dict:
    """Build standardized canonical document model"""
    
    text_blocks = result.get('text_blocks', [])
    full_text = result.get('raw_text', '')
    
    return {
        'id': request_id,
        'filename': filename,
        'file_type': file_type,
        'file_size': file_size,
        'extraction': {
            'method': result.get('method', 'unknown'),
            'confidence': result.get('confidence', 0),
            'timestamp': time.strftime('%Y-%m-%dT%H:%M:%SZ')
        },
        'structure': {
            'page_count': result.get('page_count', 1),
            'tables': result.get('tables', []),
            'images': result.get('images', [])
        },
        'content': {
            'text_blocks': text_blocks,
            'full_text': full_text
        },
        'stats': {
            'total_chars': len(full_text),
            'total_words': len(full_text.split()),
            'total_blocks': len(text_blocks),
            'total_tables': len(result.get('tables', [])),
            'total_images': len(result.get('images', []))
        }
    }

# =============================================================================
# HELPERS
# =============================================================================

def detect_file_type(filename: str, mimetype: str) -> str:
    """Detect file type from filename or mimetype"""
    
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    if ext == 'pdf' or 'pdf' in mimetype:
        return 'pdf'
    if ext in ['pptx', 'ppt'] or 'presentation' in mimetype or 'powerpoint' in mimetype:
        return 'pptx'
    if ext in ['docx', 'doc'] or 'word' in mimetype or 'document' in mimetype:
        return 'docx'
    if ext == 'txt' or 'text/plain' in mimetype:
        return 'txt'
    if ext in ['png', 'jpg', 'jpeg', 'webp', 'gif']:
        return 'image'
    
    return 'unknown'

# =============================================================================
# MAIN
# =============================================================================

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    print(f"🧠 Document Intelligence Service v2.0 starting on port {port}")
    print(f"   PDF: Document AI → PyMuPDF → OCR")
    print(f"   PPTX: python-pptx + notes + tables + image OCR")
    print(f"   DOCX: python-docx + image OCR")
    app.run(host='0.0.0.0', port=port, debug=False)
