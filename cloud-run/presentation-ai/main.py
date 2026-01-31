"""
MindSparkle AI Presentation Generator v1.0
Ultimate presentation creation with multi-AI integration

Architecture:
- GPT-4o: Slide structure, content, speaker notes
- DALL-E 3: Custom image generation
- Mermaid: Diagrams, flowcharts, charts
- Canva API: Professional templates & styles
- python-pptx: PPTX generation fallback
- ReportLab: PDF generation
"""

import os
import io
import uuid
import time
import json
import base64
import tempfile
import requests
from typing import Optional, Dict, Any, List
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# PDF generation imports
from reportlab.lib import colors as rl_colors
from reportlab.lib.pagesizes import LETTER, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Web search imports
try:
    from duckduckgo_search import DDGS
    SEARCH_ENABLED = True
except ImportError:
    SEARCH_ENABLED = False
    print('[Warning] duckduckgo-search not installed. Web search disabled.')

app = Flask(__name__)
CORS(app)

# Configuration
OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY', '')
CANVA_API_KEY = os.environ.get('CANVA_API_KEY', '')
MIDJOURNEY_API_KEY = os.environ.get('MIDJOURNEY_API_KEY', '')  # Via proxy service

# =============================================================================
# PRESENTATION STYLES
# =============================================================================

PRESENTATION_STYLES = {
    'professional': {
        'name': 'Professional',
        'description': 'Clean, corporate look with blue accent',
        'primary_color': '2563EB',  # Blue
        'secondary_color': '1E40AF',
        'accent_color': '3B82F6',
        'background': 'FFFFFF',
        'text_color': '1F2937',
        'font_title': 'Arial',
        'font_body': 'Arial',
    },
    'modern': {
        'name': 'Modern',
        'description': 'Bold, contemporary design with gradients',
        'primary_color': '7C3AED',  # Purple
        'secondary_color': '5B21B6',
        'accent_color': 'A78BFA',
        'background': 'FAFAFA',
        'text_color': '111827',
        'font_title': 'Helvetica',
        'font_body': 'Helvetica',
    },
    'minimal': {
        'name': 'Minimal',
        'description': 'Simple, elegant with lots of whitespace',
        'primary_color': '000000',
        'secondary_color': '374151',
        'accent_color': '6B7280',
        'background': 'FFFFFF',
        'text_color': '111827',
        'font_title': 'Helvetica',
        'font_body': 'Helvetica',
    },
    'creative': {
        'name': 'Creative',
        'description': 'Colorful, dynamic with bold elements',
        'primary_color': 'EC4899',  # Pink
        'secondary_color': 'F59E0B',  # Orange
        'accent_color': '10B981',  # Green
        'background': 'FFF7ED',
        'text_color': '1F2937',
        'font_title': 'Arial Black',
        'font_body': 'Arial',
    },
    'dark': {
        'name': 'Dark Mode',
        'description': 'Dark background with light text',
        'primary_color': '60A5FA',  # Light blue
        'secondary_color': '34D399',  # Green
        'accent_color': 'F472B6',  # Pink
        'background': '111827',
        'text_color': 'F9FAFB',
        'font_title': 'Arial',
        'font_body': 'Arial',
    },
    'academic': {
        'name': 'Academic',
        'description': 'Formal, scholarly presentation style',
        'primary_color': '1E3A5F',  # Navy
        'secondary_color': '7C2D12',  # Brown
        'accent_color': '047857',  # Dark green
        'background': 'FFFBEB',
        'text_color': '1F2937',
        'font_title': 'Times New Roman',
        'font_body': 'Georgia',
    },
    'startup': {
        'name': 'Startup Pitch',
        'description': 'High-energy, investor-ready design',
        'primary_color': 'EF4444',  # Red
        'secondary_color': 'F97316',  # Orange
        'accent_color': 'FBBF24',  # Yellow
        'background': 'FFFFFF',
        'text_color': '0F172A',
        'font_title': 'Arial Black',
        'font_body': 'Arial',
    },
    'education': {
        'name': 'Education',
        'description': 'Friendly, engaging for learning',
        'primary_color': '0891B2',  # Cyan
        'secondary_color': '0D9488',  # Teal
        'accent_color': 'F59E0B',  # Amber
        'background': 'F0FDFA',
        'text_color': '134E4A',
        'font_title': 'Arial',
        'font_body': 'Arial',
    },
}

# =============================================================================
# HEALTH CHECK
# =============================================================================

@app.route('/', methods=['GET'])
def index():
    return jsonify({
        'service': 'MindSparkle AI Presentation Generator',
        'version': '2.0.0',
        'status': 'healthy',
        'ai_vendors': {
            'content': 'OpenAI GPT-4o',
            'images': 'DALL-E 3 / Midjourney',
            'diagrams': 'Mermaid.js',
            'templates': 'Canva API + Custom',
            'web_search': 'DuckDuckGo' if SEARCH_ENABLED else 'Disabled'
        },
        'features': {
            'basic_generation': True,
            'enhanced_generation': SEARCH_ENABLED,
            'web_search_enrichment': SEARCH_ENABLED,
            'professional_slides': True,
            'pdf_export': True,
            'pptx_export': True
        },
        'styles': list(PRESENTATION_STYLES.keys()),
        'endpoints': {
            'POST /generate': 'Generate presentation from document',
            'POST /generate-enhanced': 'Generate with web search enrichment',
            'POST /generate-pdf': 'Generate PDF presentation',
            'POST /generate-pdf-enhanced': 'Generate enhanced PDF with web search',
            'POST /generate-slide': 'Generate single slide',
            'POST /preview': 'Preview slide structure',
            'GET /styles': 'Get available styles',
            'GET /download/:id': 'Download generated PPTX',
            'GET /download-pdf/:id': 'Download generated PDF',
        }
    })

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'timestamp': time.strftime('%Y-%m-%dT%H:%M:%SZ')})

@app.route('/styles', methods=['GET'])
def get_styles():
    return jsonify({'styles': PRESENTATION_STYLES})

@app.route('/canva/status', methods=['GET'])
def canva_status():
    """Check Canva API integration status"""
    is_configured = bool(CANVA_API_KEY)
    return jsonify({
        'canva_enabled': is_configured,
        'message': 'Canva API is configured and ready' if is_configured else 'Canva API key not configured. Using python-pptx fallback.',
        'fallback': 'python-pptx'
    })

# =============================================================================
# WEB SEARCH FOR ENHANCED PRESENTATIONS
# =============================================================================

def search_web_for_topic(query: str, max_results: int = 5) -> List[Dict[str, str]]:
    """Search the web using DuckDuckGo for relevant information"""
    
    if not SEARCH_ENABLED:
        print('[Search] Web search not available')
        return []
    
    try:
        print(f'[Search] Searching for: {query}')
        with DDGS() as ddgs:
            results = []
            for r in ddgs.text(query, max_results=max_results):
                results.append({
                    'title': r.get('title', ''),
                    'body': r.get('body', ''),
                    'url': r.get('href', '')
                })
            print(f'[Search] Found {len(results)} results')
            return results
    except Exception as e:
        print(f'[Search] Error: {e}')
        return []

def extract_search_queries_from_content(content: str, slide_count: int = 10) -> List[str]:
    """Use GPT-4o to extract relevant search queries from document content"""
    
    prompt = f"""Analyze this document and generate {min(slide_count - 2, 5)} concise search queries to find relevant data, statistics, images, and supporting information that would enhance a professional presentation.

DOCUMENT CONTENT:
{content[:8000]}

Generate search queries that will find:
- Current statistics and data
- Visual examples and case studies
- Expert opinions and research
- Relevant images and diagrams

Return ONLY a JSON array of 3-5 search query strings. Keep queries specific and relevant.
Example: ["artificial intelligence market growth 2025", "machine learning use cases healthcare", "AI statistics 2025"]"""

    try:
        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {OPENAI_API_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-4o',
                'messages': [
                    {'role': 'system', 'content': 'You are a research assistant. Extract relevant search queries. Return only valid JSON array.'},
                    {'role': 'user', 'content': prompt}
                ],
                'temperature': 0.7,
                'max_tokens': 500
            },
            timeout=30
        )
        
        result = response.json()
        if 'error' in result:
            raise Exception(f"OpenAI API error: {result['error'].get('message', 'Unknown error')}")
        
        content_text = result['choices'][0]['message']['content']
        
        # Parse JSON from response
        json_match = content_text
        if '```json' in content_text:
            json_match = content_text.split('```json')[1].split('```')[0]
        elif '```' in content_text:
            json_match = content_text.split('```')[1].split('```')[0]
        
        queries = json.loads(json_match.strip())
        print(f'[Search] Extracted {len(queries)} search queries')
        return queries if isinstance(queries, list) else []
        
    except Exception as e:
        print(f'[Search] Error extracting queries: {e}')
        return []

def generate_enhanced_slide_structure(content: str, slide_count: int = 10, style: str = 'professional', web_results: List[Dict] = None) -> List[Dict]:
    """Generate slide structure enriched with web search results"""
    
    # Prepare web search context
    web_context = ""
    if web_results and len(web_results) > 0:
        web_context = "\n\nADDITIONAL RESEARCH DATA FROM WEB:\n"
        for i, result in enumerate(web_results[:10], 1):
            web_context += f"\n{i}. {result.get('title', '')}\n{result.get('body', '')[:200]}...\nSource: {result.get('url', '')}\n"
    
    prompt = f"""Analyze this document and web research, then create a PROFESSIONAL {slide_count}-slide presentation.

DOCUMENT CONTENT:
{content[:12000]}
{web_context}

CREATE A PRESENTATION WITH:
1. **Opening Title Slide** - Professional title slide with document name
2. **Agenda/Overview** - Clear outline of presentation structure
3. **Main Content Slides** - Rich content with data from document AND web research
4. **Closing Thank You Slide** - Professional closing with key takeaway

IMPORTANT REQUIREMENTS:
- Use data, statistics, and insights from BOTH the document AND web research
- Include specific numbers, dates, and facts from web research where relevant
- Create diverse slide layouts (content, two_column, image_focus, chart, diagram)
- Each slide should have 3-5 bullet points maximum
- Include professional image prompts that match the content
- Add diagrams (flowchart, timeline, comparison) where appropriate
- Include charts (bar, pie, line) when presenting data/statistics

FOR EACH SLIDE, PROVIDE:
- slide_type: "title" | "agenda" | "content" | "two_column" | "image_focus" | "chart" | "diagram" | "summary" | "closing"
- title: Compelling, professional slide title
- subtitle: Optional subtitle (especially for title and closing slides)
- bullet_points: Array of 3-5 key points (use data from web research)
- image_prompt: Detailed DALL-E prompt for relevant professional image
- diagram_type: "flowchart" | "timeline" | "comparison" | "hierarchy" | "cycle" | null
- diagram_data: Mermaid diagram code if diagram_type is set
- chart_type: "bar" | "pie" | "line" | null
- chart_data: object with "labels" array and "values" array if chart_type is set
- speaker_notes: Detailed notes (2-3 sentences, include sources from web when relevant)
- layout: "full_image" | "left_image" | "right_image" | "top_image" | "no_image"
- web_sources: Array of relevant web source URLs used in this slide (if any)

FIRST SLIDE MUST BE:
{{
  "slide_type": "title",
  "title": "[Document Title]",
  "subtitle": "Professional Presentation",
  "layout": "no_image",
  "speaker_notes": "Welcome and introduction"
}}

LAST SLIDE MUST BE:
{{
  "slide_type": "closing",
  "title": "Thank You",
  "subtitle": "Questions & Discussion",
  "bullet_points": ["Key Takeaway 1", "Key Takeaway 2", "Key Takeaway 3"],
  "layout": "no_image",
  "speaker_notes": "Thank the audience and invite questions"
}}

Return ONLY valid JSON array. Make it PROFESSIONAL, DATA-DRIVEN, and ENGAGING using the {style} style."""

    try:
        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {OPENAI_API_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-4o',
                'messages': [
                    {'role': 'system', 'content': 'You are an expert presentation designer and researcher. Create stunning, professional, data-driven presentations. Return only valid JSON.'},
                    {'role': 'user', 'content': prompt}
                ],
                'temperature': 0.7,
                'max_tokens': 4096
            },
            timeout=90
        )
        
        result = response.json()
        
        if 'error' in result:
            raise Exception(f"OpenAI API error: {result['error'].get('message', 'Unknown error')}")
        
        content_text = result['choices'][0]['message']['content']
        print(f'[GPT-4o Enhanced] Raw response length: {len(content_text)}')
        
        # Parse JSON from response
        json_match = content_text
        if '```json' in content_text:
            json_match = content_text.split('```json')[1].split('```')[0]
        elif '```' in content_text:
            json_match = content_text.split('```')[1].split('```')[0]
        
        slides = json.loads(json_match.strip())
        print(f'[GPT-4o Enhanced] Parsed {len(slides)} slides')
        return slides
        
    except Exception as e:
        print(f'[GPT-4o Enhanced] Error: {e}')
        import traceback
        traceback.print_exc()
        
        # Fallback to basic structure
        return generate_slide_structure(content, slide_count, style)

# =============================================================================
# GPT-4O SLIDE STRUCTURE GENERATION
# =============================================================================

def generate_slide_structure(content: str, slide_count: int = 10, style: str = 'professional') -> List[Dict]:
    """Use GPT-4o to generate optimal slide structure from document content"""
    
    prompt = f"""Analyze this document and create a {slide_count}-slide presentation structure.

DOCUMENT CONTENT:
{content[:15000]}

CREATE A PRESENTATION WITH:
1. Title slide
2. Agenda/Overview slide  
3. Main content slides (with variety of layouts)
4. Summary/Conclusion slide

FOR EACH SLIDE, PROVIDE:
- slide_type: "title" | "agenda" | "content" | "two_column" | "image_focus" | "chart" | "diagram" | "quote" | "summary"
- title: Compelling slide title
- subtitle: Optional subtitle
- bullet_points: Array of key points (3-5 max)
- image_prompt: Detailed DALL-E prompt for a relevant, professional image (be specific about style, colors, composition)
- diagram_type: If diagram needed - "flowchart" | "timeline" | "comparison" | "hierarchy" | "cycle" | null
- diagram_data: Mermaid diagram code if diagram_type is set
- chart_type: If chart needed - "bar" | "pie" | "line" | null
- chart_data: Chart data if chart_type is set
- speaker_notes: Detailed speaker notes (2-3 sentences)
- layout: "full_image" | "left_image" | "right_image" | "top_image" | "no_image"

Return ONLY valid JSON array of slides. Make it PROFESSIONAL and ENGAGING.
Use the {style} style - adjust tone and formality accordingly."""

    try:
        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {OPENAI_API_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-4o',
                'messages': [
                    {'role': 'system', 'content': 'You are an expert presentation designer. Create visually stunning, professional presentations. Return only valid JSON.'},
                    {'role': 'user', 'content': prompt}
                ],
                'temperature': 0.7,
                'max_tokens': 4096
            },
            timeout=60
        )
        
        result = response.json()
        
        # Check for API errors
        if 'error' in result:
            raise Exception(f"OpenAI API error: {result['error'].get('message', 'Unknown error')}")
        
        content = result['choices'][0]['message']['content']
        print(f'[GPT-4o] Raw response length: {len(content)}')
        
        # Parse JSON from response
        json_match = content
        if '```json' in content:
            json_match = content.split('```json')[1].split('```')[0]
        elif '```' in content:
            json_match = content.split('```')[1].split('```')[0]
        
        slides = json.loads(json_match.strip())
        print(f'[GPT-4o] Parsed {len(slides)} slides')
        return slides
        
    except Exception as e:
        print(f'[GPT-4o] Error generating structure: {e}')
        import traceback
        traceback.print_exc()
        
        # Generate proper fallback structure based on slide_count
        fallback_slides = [
            {'slide_type': 'title', 'title': 'Presentation', 'subtitle': 'AI Generated', 'layout': 'no_image', 'bullet_points': [], 'speaker_notes': 'Welcome to this presentation.'},
            {'slide_type': 'agenda', 'title': 'Agenda', 'bullet_points': ['Introduction', 'Key Points', 'Summary'], 'layout': 'no_image', 'speaker_notes': 'Overview of topics.'},
        ]
        
        # Add content slides
        content_snippets = content[:3000].split('\n\n')[:slide_count - 3]
        for i, snippet in enumerate(content_snippets):
            if snippet.strip():
                fallback_slides.append({
                    'slide_type': 'content',
                    'title': f'Key Point {i + 1}',
                    'bullet_points': [s.strip()[:200] for s in snippet.split('. ')[:4] if s.strip()],
                    'layout': 'no_image',
                    'speaker_notes': 'Key information from the document.'
                })
        
        # Add summary slide
        fallback_slides.append({
            'slide_type': 'summary',
            'title': 'Summary',
            'bullet_points': ['Key takeaways from this presentation'],
            'layout': 'no_image',
            'speaker_notes': 'Thank you for your attention.'
        })
        
        return fallback_slides[:slide_count]


def generate_topic_slide_structure(content: str, max_topics: int = 10, style: str = 'professional') -> List[Dict]:
    """Generate a deck where each document topic becomes a slide heading.

    Product rules:
    - No title/agenda/closing slides
    - No author/presenter/company content
    - No speaker notes
    - Each slide includes: image + blocks + text, plus optional table/diagram when relevant
    """

    prompt = f"""You are generating a professional slide deck from a document.

DOCUMENT CONTENT:
{content[:15000]}

GOAL:
- Extract the most important topics from the document.
- Each topic becomes ONE slide.

HARD RULES:
- Do NOT create a title slide, agenda slide, summary slide, or closing slide.
- Do NOT include any author name, presenter name, company name, or speaker notes.
- Do NOT include any text like "AI Generated".
- Slide titles must be the extracted topics.

FOR EACH SLIDE RETURN THESE FIELDS:
- slide_type: must be "topic"
- title: topic heading
- overview: 1–2 sentence professional explanation
- blocks: 2–3 short labels (3–6 words) for key sub-ideas
- bullet_points: 3–5 concise bullets
- image_prompt: a DALL-E prompt for a clean educational visual (NO TEXT)
- table: optional object with title, headers, rows (use when comparison/structured data fits)
- diagram_type: optional one of "flowchart"|"timeline"|"comparison"|"hierarchy"|"cycle"|null
- diagram_data: Mermaid code when diagram_type is set (NO markdown fences)
- layout: always "right_image"

QUALITY RULES:
- Ensure the overall deck contains at least ONE slide with a table and at least ONE slide with a diagram when the content allows.
- Avoid fluff; be accurate to the document.

Return ONLY a valid JSON array with up to {max_topics} slides.
Use the {style} tone/style."""

    try:
        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {OPENAI_API_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-4o',
                'messages': [
                    {'role': 'system', 'content': 'You are an expert presentation designer. Return only valid JSON.'},
                    {'role': 'user', 'content': prompt}
                ],
                'temperature': 0.5,
                'max_tokens': 4096
            },
            timeout=75
        )

        result = response.json()
        if 'error' in result:
            raise Exception(f"OpenAI API error: {result['error'].get('message', 'Unknown error')}")

        content_text = result['choices'][0]['message']['content']
        json_match = content_text
        if '```json' in content_text:
            json_match = content_text.split('```json')[1].split('```')[0]
        elif '```' in content_text:
            json_match = content_text.split('```')[1].split('```')[0]

        slides = json.loads(json_match.strip())
        if not isinstance(slides, list):
            raise Exception('Invalid JSON response (not an array)')

        normalized: List[Dict[str, Any]] = []
        for s in slides[:max_topics]:
            if not isinstance(s, dict):
                continue
            normalized.append({
                'slide_type': 'topic',
                'title': str(s.get('title', '')).strip(),
                'overview': str(s.get('overview', '')).strip(),
                'blocks': s.get('blocks', []) if isinstance(s.get('blocks'), list) else [],
                'bullet_points': s.get('bullet_points', []) if isinstance(s.get('bullet_points'), list) else [],
                'image_prompt': str(s.get('image_prompt', '')).strip(),
                'table': s.get('table') if isinstance(s.get('table'), dict) else None,
                'diagram_type': s.get('diagram_type'),
                'diagram_data': str(s.get('diagram_data', '')).strip() if s.get('diagram_data') else None,
                'layout': 'right_image',
            })

        if len(normalized) == 0:
            raise Exception('No slides returned')

        return normalized

    except Exception as e:
        print(f'[GPT-4o Topics] Error generating topic structure: {e}')
        return [
            {
                'slide_type': 'topic',
                'title': 'Key Topics',
                'overview': 'Key topics extracted from the document.',
                'blocks': ['Overview', 'Concepts', 'Summary'],
                'bullet_points': ['Topic extraction failed. Please try again.'],
                'image_prompt': 'Abstract educational illustration, clean, modern, no text',
                'table': None,
                'diagram_type': None,
                'diagram_data': None,
                'layout': 'right_image',
            }
        ]

# =============================================================================
# DALL-E 3 IMAGE GENERATION
# =============================================================================

def generate_image_dalle(prompt: str, style: str = 'professional', size: str = '1792x1024') -> Optional[str]:
    """Generate image using DALL-E 3"""
    
    style_modifiers = {
        'professional': 'professional corporate style, clean modern design, subtle colors',
        'modern': 'modern minimalist style, bold geometric shapes, vibrant gradients',
        'minimal': 'minimal clean style, lots of white space, simple elegant',
        'creative': 'creative colorful style, dynamic composition, artistic flair',
        'dark': 'dark moody style, dramatic lighting, sleek modern',
        'academic': 'academic scholarly style, classic elegant, sophisticated',
        'startup': 'startup tech style, innovative dynamic, energetic modern',
        'education': 'educational friendly style, warm inviting, clear informative',
    }
    
    full_prompt = f"{prompt}. Style: {style_modifiers.get(style, style_modifiers['professional'])}. High quality, 4K, professional presentation slide image."
    
    try:
        response = requests.post(
            'https://api.openai.com/v1/images/generations',
            headers={
                'Authorization': f'Bearer {OPENAI_API_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'dall-e-3',
                'prompt': full_prompt,
                'n': 1,
                'size': size,
                'quality': 'hd',
                'style': 'vivid' if style in ['creative', 'startup'] else 'natural'
            },
            timeout=60
        )
        
        result = response.json()
        if 'data' in result and len(result['data']) > 0:
            image_url = result['data'][0]['url']
            
            # Download and convert to base64
            img_response = requests.get(image_url, timeout=30)
            if img_response.status_code == 200:
                return base64.b64encode(img_response.content).decode('utf-8')
        
        return None
        
    except Exception as e:
        print(f'[DALL-E 3] Error generating image: {e}')
        return None


def generate_image_nano_banana(prompt: str, style: str = 'professional', size: str = '1024x1024') -> Optional[str]:
    """Generate image using OpenAI gpt-image-1 (labeled as 'Nano Banana' in the app)."""

    # Reuse the same style modifiers as DALL·E for consistent aesthetics.
    style_modifiers = {
        'professional': 'professional corporate style, clean modern design, subtle colors',
        'modern': 'modern minimalist style, bold geometric shapes, vibrant gradients',
        'minimal': 'minimal clean style, lots of white space, simple elegant',
        'creative': 'creative colorful style, dynamic composition, artistic flair',
        'dark': 'dark moody style, dramatic lighting, sleek modern',
        'academic': 'academic scholarly style, classic elegant, sophisticated',
        'startup': 'startup tech style, innovative dynamic, energetic modern',
        'education': 'educational friendly style, warm inviting, clear informative',
    }

    full_prompt = f"{prompt}. Style: {style_modifiers.get(style, style_modifiers['professional'])}. High quality, realistic, professional presentation slide image."

    try:
        response = requests.post(
            'https://api.openai.com/v1/images/generations',
            headers={
                'Authorization': f'Bearer {OPENAI_API_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-image-1',
                'prompt': full_prompt,
                'n': 1,
                'size': size,
            },
            timeout=90
        )

        # If the account doesn't have gpt-image-1 access, fall back gracefully.
        if response.status_code in (400, 401, 403):
            print(f'[Nano Banana] gpt-image-1 not available (status={response.status_code}); falling back to DALL-E 3')
            return generate_image_dalle(prompt, style, size='1792x1024')

        result = response.json()
        if 'data' in result and len(result['data']) > 0:
            # Prefer base64 if present.
            b64 = result['data'][0].get('b64_json')
            if b64:
                return b64

            image_url = result['data'][0].get('url')
            if image_url:
                img_response = requests.get(image_url, timeout=30)
                if img_response.status_code == 200:
                    return base64.b64encode(img_response.content).decode('utf-8')

        return None
    except Exception as e:
        print(f'[Nano Banana] Error generating image: {e}')
        return None


def generate_image_midjourney(prompt: str, style: str = 'professional') -> Optional[str]:
    """Generate image using Midjourney (optional).

    This code path is only used when MIDJOURNEY_API_KEY and MIDJOURNEY_API_URL are configured.
    If not configured (or the call fails), callers should fall back to DALL·E.
    """

    api_key = (MIDJOURNEY_API_KEY or '').strip()
    api_url = (os.environ.get('MIDJOURNEY_API_URL', '') or '').strip()
    if not api_key or not api_url:
        return None

    try:
        resp = requests.post(
            api_url,
            headers={
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            },
            json={
                'prompt': prompt,
                'style': style,
            },
            timeout=120
        )
        if resp.status_code != 200:
            print(f'[Midjourney] Non-OK status: {resp.status_code}')
            return None

        data = resp.json() if resp.text else {}

        # Support either base64 or url payloads.
        b64 = data.get('b64') or data.get('base64') or data.get('image_base64')
        if isinstance(b64, str) and b64.strip():
            return b64.strip()

        url = data.get('url') or data.get('image_url')
        if isinstance(url, str) and url.strip():
            img_response = requests.get(url.strip(), timeout=60)
            if img_response.status_code == 200:
                return base64.b64encode(img_response.content).decode('utf-8')

        return None
    except Exception as e:
        print(f'[Midjourney] Error generating image: {e}')
        return None

# =============================================================================
# MERMAID DIAGRAM GENERATION
# =============================================================================

def generate_mermaid_diagram(diagram_type: str, diagram_data: str) -> Optional[str]:
    """Generate diagram image from Mermaid code"""
    
    try:
        # Use Mermaid.ink API to render diagram
        mermaid_code = diagram_data.strip()
        encoded = base64.b64encode(mermaid_code.encode('utf-8')).decode('utf-8')
        
        # Mermaid.ink API
        url = f'https://mermaid.ink/img/{encoded}?type=png&bgColor=transparent'
        
        response = requests.get(url, timeout=30)
        if response.status_code == 200:
            return base64.b64encode(response.content).decode('utf-8')
        
        return None
        
    except Exception as e:
        print(f'[Mermaid] Error generating diagram: {e}')
        return None

# =============================================================================
# PPTX GENERATION
# =============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def create_presentation(
    slides: List[Dict],
    style_name: str = 'professional',
    include_images: bool = True,
    include_notes: bool = True,
    image_mode: str = 'default'
) -> bytes:
    """Create PPTX from slide structure"""
    
    style = PRESENTATION_STYLES.get(style_name, PRESENTATION_STYLES['professional'])
    
    prs = Presentation()
    # Avoid unprofessional metadata like an 'Author' showing up in PPT properties
    try:
        prs.core_properties.author = ""
        prs.core_properties.last_modified_by = ""
        prs.core_properties.company = ""
        prs.core_properties.comments = ""
    except Exception:
        pass
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    for i, slide_data in enumerate(slides):
        print(f'[PPTX] Creating slide {i+1}: {slide_data.get("slide_type", "content")}')
        
        slide_type = slide_data.get('slide_type', 'content')
        layout = slide_data.get('layout', 'no_image')
        
        # Use blank layout for full control
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(style['background'])
        
        # Handle different slide types
        if slide_type == 'title':
            create_title_slide(slide, slide_data, style)
        elif slide_type == 'closing':
            create_closing_slide(slide, slide_data, style)
        elif slide_type == 'agenda':
            create_agenda_slide(slide, slide_data, style)
        elif slide_type == 'two_column':
            create_two_column_slide(slide, slide_data, style)
        elif slide_type == 'quote':
            create_quote_slide(slide, slide_data, style)
        elif slide_type == 'summary':
            create_summary_slide(slide, slide_data, style)
        elif slide_type == 'topic':
            create_topic_slide(slide, slide_data, style)
        else:
            create_content_slide(slide, slide_data, style, layout)
        
        # Generate and add image if needed
        if include_images and slide_data.get('image_prompt') and layout != 'no_image':
            mode = (image_mode or 'default').strip().lower()
            image_base64 = None

            if mode in ('realism', 'enhance_realism', 'enhanced_realism', 'nano_banana', 'nanobanana'):
                image_base64 = generate_image_nano_banana(slide_data['image_prompt'], style_name)
            elif mode in ('premium', 'premium_visuals', 'midjourney'):
                # Optional: only works if env vars are configured.
                image_base64 = generate_image_midjourney(slide_data['image_prompt'], style_name)
                if not image_base64:
                    image_base64 = generate_image_dalle(slide_data['image_prompt'], style_name)
            else:
                # Default
                image_base64 = generate_image_dalle(slide_data['image_prompt'], style_name)

            if image_base64:
                add_image_to_slide(slide, image_base64, layout)
        
        # Generate and add diagram if needed
        if slide_data.get('diagram_type') and slide_data.get('diagram_data'):
            diagram_base64 = generate_mermaid_diagram(
                slide_data['diagram_type'],
                slide_data['diagram_data']
            )
            if diagram_base64:
                add_diagram_to_slide(slide, diagram_base64)
        
        # Add speaker notes (optional)
        if include_notes and slide_data.get('speaker_notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['speaker_notes']
    
    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes.getvalue()

def create_title_slide(slide, data: Dict, style: Dict):
    """Create title slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Presentation')
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.font.size = Pt(28)
        p.font.color.rgb = hex_to_rgb(style['secondary_color'])
        p.alignment = PP_ALIGN.CENTER

def create_closing_slide(slide, data: Dict, style: Dict):
    """Create professional closing/thank you slide"""
    # Title (Thank You)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Thank You')
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(style['secondary_color'])
        p.alignment = PP_ALIGN.CENTER
    
    # Key takeaways or closing message (if provided)
    if data.get('bullet_points') and len(data['bullet_points']) > 0:
        content_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.333), Inches(1.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(data['bullet_points'][:3]):  # Max 3 points
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(style['text_color'])
            p.space_after = Pt(8)
            p.alignment = PP_ALIGN.CENTER

def create_agenda_slide(slide, data: Dict, style: Dict):
    """Create agenda/overview slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Agenda')
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    
    # Bullet points
    if data.get('bullet_points'):
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(data['bullet_points']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(style['text_color'])
            p.space_after = Pt(18)

def create_content_slide(slide, data: Dict, style: Dict, layout: str):
    """Create standard content slide"""
    # Calculate text area based on layout
    if layout in ['left_image', 'right_image']:
        text_left = Inches(7) if layout == 'left_image' else Inches(0.5)
        text_width = Inches(5.5)
    else:
        text_left = Inches(0.5)
        text_width = Inches(12.333)
    
    # Title
    title_box = slide.shapes.add_textbox(text_left, Inches(0.5), text_width, Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    
    # Content
    if data.get('bullet_points'):
        content_box = slide.shapes.add_textbox(text_left, Inches(1.8), text_width, Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(data['bullet_points']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(style['text_color'])
            p.space_after = Pt(12)


def create_topic_slide(slide, data: Dict, style: Dict):
    """Create a topic slide: title + blocks + overview + bullets + (optional) table.

    Images are handled by the main loop via image_prompt and layout.
    Diagrams are handled by the main loop via diagram_type/diagram_data.
    """

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])

    # Blocks (pills)
    blocks = data.get('blocks', []) if isinstance(data.get('blocks'), list) else []
    blocks = [str(b).strip() for b in blocks if str(b).strip()][:3]
    x = Inches(0.7)
    y = Inches(1.25)
    for b in blocks:
        w = Inches(2.8)
        h = Inches(0.45)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(style['accent_color'])
        shape.line.color.rgb = hex_to_rgb(style['accent_color'])
        tfb = shape.text_frame
        tfb.clear()
        para = tfb.paragraphs[0]
        para.text = b
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
        x = x + Inches(3.0)

    # Left content (overview + bullets)
    left_x = Inches(0.7)
    left_w = Inches(6.1)

    overview = str(data.get('overview', '')).strip()
    if overview:
        overview_box = slide.shapes.add_textbox(left_x, Inches(1.85), left_w, Inches(1.05))
        tf = overview_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = overview
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(style['text_color'])

    bullets = data.get('bullet_points', []) if isinstance(data.get('bullet_points'), list) else []
    bullets = [str(b).strip() for b in bullets if str(b).strip()][:5]
    if bullets:
        bullets_box = slide.shapes.add_textbox(left_x, Inches(2.95), left_w, Inches(2.55))
        tf = bullets_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"• {point}"
            para.font.size = Pt(18)
            para.font.color.rgb = hex_to_rgb(style['text_color'])
            para.space_after = Pt(10)

    # Optional table
    table = data.get('table') if isinstance(data.get('table'), dict) else None
    if table:
        headers = table.get('headers', []) if isinstance(table.get('headers'), list) else []
        rows = table.get('rows', []) if isinstance(table.get('rows'), list) else []
        headers = [str(h) for h in headers][:5]
        rows = [r for r in rows if isinstance(r, list)][:6]
        if headers and rows:
            add_table_to_slide(slide, headers, rows, Inches(0.7), Inches(5.6), Inches(6.1), Inches(1.75), style)


def add_table_to_slide(slide, headers: List[str], rows: List[List[str]], x, y, w, h, style: Dict):
    """Add a compact table with basic styling."""

    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(style['primary_color'])

    for ri, row in enumerate(rows, start=1):
        for ci in range(n_cols):
            value = ''
            if ci < len(row):
                value = str(row[ci])
            cell = table.cell(ri, ci)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = hex_to_rgb(style['text_color'])
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb('F3F4F6')

def create_two_column_slide(slide, data: Dict, style: Dict):
    """Create two-column comparison slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    bullet_points = data.get('bullet_points', [])
    half = len(bullet_points) // 2
    
    for i, point in enumerate(bullet_points[:half]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(style['text_color'])
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points[half:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(style['text_color'])
        p.space_after = Pt(10)

def create_quote_slide(slide, data: Dict, style: Dict):
    """Create quote/highlight slide"""
    # Large quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    
    bullet_points = data.get('bullet_points', [''])
    p.text = f'"{bullet_points[0]}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    p.alignment = PP_ALIGN.CENTER

def create_summary_slide(slide, data: Dict, style: Dict):
    """Create summary/conclusion slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Key Takeaways')
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style['primary_color'])
    p.alignment = PP_ALIGN.CENTER
    
    # Summary points
    if data.get('bullet_points'):
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(data['bullet_points']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {point}"
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(style['accent_color'])
            p.space_after = Pt(14)
            p.alignment = PP_ALIGN.CENTER

def add_image_to_slide(slide, image_base64: str, layout: str):
    """Add image to slide based on layout"""
    try:
        image_bytes = base64.b64decode(image_base64)
        image_stream = io.BytesIO(image_bytes)
        
        if layout == 'full_image':
            slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        elif layout == 'left_image':
            slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1.5), width=Inches(6), height=Inches(5.5))
        elif layout == 'right_image':
            slide.shapes.add_picture(image_stream, Inches(6.833), Inches(1.5), width=Inches(6), height=Inches(5.5))
        elif layout == 'top_image':
            slide.shapes.add_picture(image_stream, Inches(2), Inches(1.5), width=Inches(9.333), height=Inches(3))
    except Exception as e:
        print(f'[PPTX] Error adding image: {e}')

def add_diagram_to_slide(slide, diagram_base64: str):
    """Add diagram to slide"""
    try:
        image_bytes = base64.b64decode(diagram_base64)
        image_stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(image_stream, Inches(3), Inches(2.5), width=Inches(7), height=Inches(4))
    except Exception as e:
        print(f'[PPTX] Error adding diagram: {e}')

# =============================================================================
# CANVA API INTEGRATION
# =============================================================================

# Canva presentation template IDs for different styles
CANVA_TEMPLATES = {
    'professional': 'DAGcHnE_R4M',  # Professional business template
    'modern': 'DAGcNvLQS8Q',       # Modern gradient template
    'minimal': 'DAGcK2XcP9w',      # Minimal clean template
    'creative': 'DAGcOxWkL2g',     # Creative colorful template
    'dark': 'DAGcP1YmM3k',         # Dark mode template
    'academic': 'DAGcQ3ZnN4l',     # Academic/scholarly template
    'startup': 'DAGcR5apO5m',      # Startup pitch template
    'education': 'DAGcS7bqP6n',    # Education friendly template
}

def get_canva_access_token() -> Optional[str]:
    """Get Canva access token using OAuth2"""
    if not CANVA_API_KEY:
        return None
    # For direct API key usage, return the key
    # For OAuth flow, this would handle token refresh
    return CANVA_API_KEY

def create_canva_design(title: str, style_name: str) -> Optional[Dict]:
    """Create a new Canva design from template"""
    
    token = get_canva_access_token()
    if not token:
        return None
    
    template_id = CANVA_TEMPLATES.get(style_name, CANVA_TEMPLATES['professional'])
    
    try:
        # Create design from template
        response = requests.post(
            'https://api.canva.com/rest/v1/designs',
            headers={
                'Authorization': f'Bearer {token}',
                'Content-Type': 'application/json'
            },
            json={
                'design_type': 'presentation',
                'title': title,
                'template_id': template_id
            },
            timeout=30
        )
        
        if response.status_code == 200:
            return response.json()
        else:
            print(f'[Canva] Create design failed: {response.status_code} - {response.text}')
            return None
            
    except Exception as e:
        print(f'[Canva] Error creating design: {e}')
        return None

def add_canva_page(design_id: str, slide_data: Dict) -> bool:
    """Add a page to Canva design"""
    
    token = get_canva_access_token()
    if not token:
        return False
    
    try:
        # Add page with content
        response = requests.post(
            f'https://api.canva.com/rest/v1/designs/{design_id}/pages',
            headers={
                'Authorization': f'Bearer {token}',
                'Content-Type': 'application/json'
            },
            json={
                'title': slide_data.get('title', ''),
                'elements': build_canva_elements(slide_data)
            },
            timeout=30
        )
        
        return response.status_code == 200
        
    except Exception as e:
        print(f'[Canva] Error adding page: {e}')
        return False

def build_canva_elements(slide_data: Dict) -> List[Dict]:
    """Build Canva elements from slide data"""
    elements = []
    
    # Title element
    if slide_data.get('title'):
        elements.append({
            'type': 'TEXT',
            'content': slide_data['title'],
            'style': {
                'font_family': 'Montserrat',
                'font_size': 48,
                'font_weight': 'bold',
                'text_align': 'center'
            },
            'position': {'x': 50, 'y': 50, 'width': 700, 'height': 80}
        })
    
    # Subtitle element
    if slide_data.get('subtitle'):
        elements.append({
            'type': 'TEXT',
            'content': slide_data['subtitle'],
            'style': {
                'font_family': 'Open Sans',
                'font_size': 24,
                'text_align': 'center'
            },
            'position': {'x': 50, 'y': 140, 'width': 700, 'height': 40}
        })
    
    # Bullet points
    if slide_data.get('bullet_points'):
        bullet_text = '\n'.join([f'• {point}' for point in slide_data['bullet_points']])
        elements.append({
            'type': 'TEXT',
            'content': bullet_text,
            'style': {
                'font_family': 'Open Sans',
                'font_size': 18,
                'line_height': 1.5
            },
            'position': {'x': 50, 'y': 200, 'width': 700, 'height': 300}
        })
    
    return elements

def export_canva_design(design_id: str, format: str = 'pptx') -> Optional[str]:
    """Export Canva design to file"""
    
    token = get_canva_access_token()
    if not token:
        return None
    
    try:
        # Start export job
        response = requests.post(
            f'https://api.canva.com/rest/v1/designs/{design_id}/exports',
            headers={
                'Authorization': f'Bearer {token}',
                'Content-Type': 'application/json'
            },
            json={
                'format': format.upper(),
                'quality': 'professional'
            },
            timeout=30
        )
        
        if response.status_code != 200:
            print(f'[Canva] Export failed: {response.status_code}')
            return None
        
        export_data = response.json()
        export_id = export_data.get('id')
        
        # Poll for export completion
        for _ in range(30):  # Max 30 seconds
            time.sleep(1)
            status_response = requests.get(
                f'https://api.canva.com/rest/v1/exports/{export_id}',
                headers={'Authorization': f'Bearer {token}'},
                timeout=30
            )
            
            if status_response.status_code == 200:
                status_data = status_response.json()
                if status_data.get('status') == 'completed':
                    return status_data.get('urls', [{}])[0].get('url')
                elif status_data.get('status') == 'failed':
                    print(f'[Canva] Export failed: {status_data}')
                    return None
        
        return None
        
    except Exception as e:
        print(f'[Canva] Error exporting: {e}')
        return None

def create_with_canva(slides: List[Dict], style_name: str, title: str = 'AI Presentation') -> Optional[str]:
    """Create presentation using Canva API (if available)"""
    
    if not CANVA_API_KEY:
        print('[Canva] No API key configured, skipping Canva integration')
        return None
    
    try:
        print(f'[Canva] Creating presentation with {len(slides)} slides, style: {style_name}')
        
        # Step 1: Create design from template
        design = create_canva_design(title, style_name)
        if not design:
            print('[Canva] Failed to create design, falling back to PPTX')
            return None
        
        design_id = design.get('design', {}).get('id')
        if not design_id:
            print('[Canva] No design ID returned')
            return None
        
        print(f'[Canva] Created design: {design_id}')
        
        # Step 2: Add pages for each slide
        for i, slide in enumerate(slides):
            print(f'[Canva] Adding slide {i+1}/{len(slides)}')
            if not add_canva_page(design_id, slide):
                print(f'[Canva] Failed to add slide {i+1}')
        
        # Step 3: Export to PPTX
        print('[Canva] Exporting design to PPTX...')
        download_url = export_canva_design(design_id, 'pptx')
        
        if download_url:
            print(f'[Canva] Export successful: {download_url}')
            return download_url
        else:
            print('[Canva] Export failed, falling back to PPTX')
            return None
            
    except Exception as e:
        print(f'[Canva] Error: {e}')
        return None

# =============================================================================
# MAIN ENDPOINTS
# =============================================================================

@app.route('/generate', methods=['POST'])
def generate_presentation():
    """Generate complete presentation from document content"""
    
    start_time = time.time()
    request_id = str(uuid.uuid4())[:8]
    
    try:
        data = request.json
        content = data.get('content', '')
        style = data.get('style', 'professional')
        slide_count = data.get('slide_count', data.get('slideCount', 10))  # Support both naming conventions
        include_images = data.get('include_images', data.get('includeImages', True))
        image_mode = data.get('image_mode', data.get('imageMode', 'default'))
        title = data.get('title', 'AI Generated Presentation')

        structure_mode = data.get('structure', data.get('structureMode', 'classic'))
        structure_mode = (structure_mode or 'classic').strip().lower()
        topics_mode = structure_mode in ('topics', 'topic')
        include_notes = not topics_mode
        
        print(f'[{request_id}] Request params: slide_count={slide_count}, style={style}, include_images={include_images}, image_mode={image_mode}, structure={structure_mode}, content_length={len(content)}')
        
        if not content or len(content) < 100:
            return jsonify({'error': 'Content too short'}), 400
        
        print(f'[{request_id}] Generating presentation: {slide_count} slides, style: {style}, structure: {structure_mode}')
        
        # Step 1: Generate slide structure with GPT-4o
        print(f'[{request_id}] Step 1: Generating slide structure...')
        try:
            slide_count_int = int(slide_count)
        except Exception:
            slide_count_int = 10

        if topics_mode:
            slides = generate_topic_slide_structure(content, slide_count_int, style)
            for s in slides:
                if isinstance(s, dict):
                    s.pop('speaker_notes', None)
        else:
            slides = generate_slide_structure(content, slide_count_int, style)
            # Ensure title slide has correct title
            if slides and slides[0].get('slide_type') == 'title':
                slides[0]['title'] = title
        
        # Step 2: Try Canva first, fall back to python-pptx
        print(f'[{request_id}] Step 2: Creating presentation...')
        canva_url = None
        if not topics_mode:
            canva_url = create_with_canva(slides, style, title)
        
        if canva_url:
            return jsonify({
                'success': True,
                'request_id': request_id,
                'download_url': canva_url,
                'provider': 'canva',
                'slide_count': len(slides),
                'style': style,
                'duration': round(time.time() - start_time, 2)
            })
        
        # Generate with python-pptx
        print(f'[{request_id}] Step 3: Building PPTX with images...')
        pptx_bytes = create_presentation(slides, style, include_images, include_notes=include_notes, image_mode=image_mode)
        
        # Save to temp file
        filename = f'presentation_{request_id}.pptx'
        filepath = os.path.join(tempfile.gettempdir(), filename)
        with open(filepath, 'wb') as f:
            f.write(pptx_bytes)
        
        duration = round(time.time() - start_time, 2)
        print(f'[{request_id}] Presentation created in {duration}s')
        
        return jsonify({
            'success': True,
            'request_id': request_id,
            'download_url': f'/download/{request_id}',
            'filename': filename,
            'provider': 'pptx',
            'slide_count': len(slides),
            'style': style,
            'structure': structure_mode,
            'include_notes': include_notes,
            'slides': slides,  # Include structure for preview
            'duration': duration
        })
        
    except Exception as e:
        print(f'[{request_id}] Error: {e}')
        return jsonify({'error': str(e), 'request_id': request_id}), 500

@app.route('/generate-enhanced', methods=['POST'])
def generate_enhanced_presentation():
    """Generate enhanced presentation with web search enrichment"""
    
    start_time = time.time()
    request_id = str(uuid.uuid4())[:8]
    
    try:
        data = request.json
        content = data.get('content', '')
        style = data.get('style', 'professional')
        slide_count = data.get('slide_count', data.get('slideCount', 10))
        include_images = data.get('include_images', data.get('includeImages', True))
        image_mode = data.get('image_mode', data.get('imageMode', 'default'))
        title = data.get('title', 'AI Generated Presentation')

        structure_mode = data.get('structure', data.get('structureMode', 'classic'))
        structure_mode = (structure_mode or 'classic').strip().lower()
        topics_mode = structure_mode in ('topics', 'topic')
        include_notes = not topics_mode

        try:
            slide_count_int = int(slide_count)
        except Exception:
            slide_count_int = 10
        
        print(f'[{request_id}] ENHANCED Request: slide_count={slide_count}, style={style}, structure={structure_mode}, content_length={len(content)}')
        
        if not content or len(content) < 100:
            return jsonify({'error': 'Content too short'}), 400
        
        all_web_results = []
        search_queries = []

        if topics_mode:
            print(f'[{request_id}] Topics mode: skipping web enrichment; using topic structure')
            slides = generate_topic_slide_structure(content, slide_count_int, style)
            for s in slides:
                if isinstance(s, dict):
                    s.pop('speaker_notes', None)
        else:
            # Step 1: Extract search queries from content
            print(f'[{request_id}] Step 1: Extracting search queries...')
            search_queries = extract_search_queries_from_content(content, slide_count_int)

            # Step 2: Perform web searches
            print(f'[{request_id}] Step 2: Searching web for enrichment data...')
            for query in search_queries[:5]:  # Limit to 5 queries
                results = search_web_for_topic(query, max_results=3)
                all_web_results.extend(results)
                time.sleep(0.5)  # Rate limiting

            print(f'[{request_id}] Collected {len(all_web_results)} web results')

            # Step 3: Generate enhanced slide structure with web data
            print(f'[{request_id}] Step 3: Generating enhanced slide structure...')
            slides = generate_enhanced_slide_structure(content, slide_count_int, style, all_web_results)

            # Ensure professional opening and closing
            if slides and slides[0].get('slide_type') == 'title':
                slides[0]['title'] = title
                slides[0]['subtitle'] = 'Professional Presentation'

            if slides and slides[-1].get('slide_type') != 'closing':
                # Add closing slide if not present
                slides.append({
                    'slide_type': 'closing',
                    'title': 'Thank You',
                    'subtitle': 'Questions & Discussion',
                    'bullet_points': ['Thank you for your attention', 'Questions welcome'],
                    'layout': 'no_image',
                    'speaker_notes': 'Thank the audience and invite questions'
                })
        
        # Step 4: Try Canva first, fall back to python-pptx
        print(f'[{request_id}] Step 4: Creating presentation...')
        canva_url = None
        if not topics_mode:
            canva_url = create_with_canva(slides, style, title)
        
        if canva_url:
            return jsonify({
                'success': True,
                'request_id': request_id,
                'download_url': canva_url,
                'provider': 'canva',
                'slide_count': len(slides),
                'style': style,
                'enhanced': True,
                'web_results_count': len(all_web_results),
                'duration': round(time.time() - start_time, 2)
            })
        
        # Generate with python-pptx
        print(f'[{request_id}] Step 5: Building PPTX with images...')
        pptx_bytes = create_presentation(slides, style, include_images, include_notes=include_notes, image_mode=image_mode)
        
        # Save to temp file
        filename = f'presentation_{request_id}.pptx'
        filepath = os.path.join(tempfile.gettempdir(), filename)
        with open(filepath, 'wb') as f:
            f.write(pptx_bytes)
        
        duration = round(time.time() - start_time, 2)
        print(f'[{request_id}] Enhanced presentation created in {duration}s')
        
        return jsonify({
            'success': True,
            'request_id': request_id,
            'download_url': f'/download/{request_id}',
            'filename': filename,
            'provider': 'pptx',
            'slide_count': len(slides),
            'style': style,
            'slides': slides,
            'structure': structure_mode,
            'include_notes': include_notes,
            'enhanced': (not topics_mode),
            'web_results_count': len(all_web_results),
            'search_queries': search_queries,
            'duration': duration
        })
        
    except Exception as e:
        print(f'[{request_id}] Enhanced Error: {e}')
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e), 'request_id': request_id}), 500

@app.route('/download/<request_id>', methods=['GET'])
def download_presentation(request_id: str):
    """Download generated presentation"""
    
    filename = f'presentation_{request_id}.pptx'
    filepath = os.path.join(tempfile.gettempdir(), filename)
    
    if not os.path.exists(filepath):
        return jsonify({'error': 'Presentation not found or expired'}), 404
    
    return send_file(
        filepath,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name=filename
    )

@app.route('/preview', methods=['POST'])
def preview_slides():
    """Generate slide structure preview without images (faster)"""
    
    try:
        data = request.json
        content = data.get('content', '')
        style = data.get('style', 'professional')
        slide_count = data.get('slide_count', data.get('slideCount', 10))  # Support both naming conventions

        structure_mode = data.get('structure', data.get('structureMode', 'classic'))
        structure_mode = (structure_mode or 'classic').strip().lower()
        topics_mode = structure_mode in ('topics', 'topic')
        
        print(f'[Preview] Request: slide_count={slide_count}, style={style}, structure={structure_mode}, content_length={len(content)}')
        
        if not content or len(content) < 100:
            return jsonify({'error': 'Content too short'}), 400
        
        try:
            slide_count_int = int(slide_count)
        except Exception:
            slide_count_int = 10

        if topics_mode:
            slides = generate_topic_slide_structure(content, slide_count_int, style)
            for s in slides:
                if isinstance(s, dict):
                    s.pop('speaker_notes', None)
        else:
            slides = generate_slide_structure(content, slide_count_int, style)
        print(f'[Preview] Generated {len(slides)} slides')
        
        return jsonify({
            'success': True,
            'slides': slides,
            'style': style,
            'structure': structure_mode,
            'style_details': PRESENTATION_STYLES.get(style)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# =============================================================================
# PDF GENERATION
# =============================================================================

def create_pdf_presentation(slides: List[Dict], style_name: str = 'professional', include_notes: bool = True) -> bytes:
    """Create PDF presentation from slide structure"""
    
    style = PRESENTATION_STYLES.get(style_name, PRESENTATION_STYLES['professional'])
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(LETTER),
        rightMargin=0.5*inch,
        leftMargin=0.5*inch,
        topMargin=0.5*inch,
        bottomMargin=0.5*inch
    )
    
    # Custom styles
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontSize=32,
        textColor=rl_colors.HexColor('#' + style['primary_color']),
        alignment=TA_CENTER,
        spaceAfter=20
    )
    
    subtitle_style = ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=rl_colors.HexColor('#' + style['secondary_color']),
        alignment=TA_CENTER,
        spaceAfter=30
    )
    
    bullet_style = ParagraphStyle(
        'SlideBullet',
        parent=styles['Normal'],
        fontSize=16,
        textColor=rl_colors.HexColor('#' + style['text_color']),
        leftIndent=40,
        spaceAfter=12,
        bulletIndent=20
    )
    
    elements = []
    
    for i, slide in enumerate(slides):
        # Add page break between slides (except first)
        if i > 0:
            elements.append(PageBreak())
        
        # Slide number
        slide_num = Paragraph(f"Slide {i+1}", ParagraphStyle('SlideNum', fontSize=10, textColor=rl_colors.gray))
        elements.append(slide_num)
        elements.append(Spacer(1, 20))
        
        # Title
        title = slide.get('title', f'Slide {i+1}')
        elements.append(Paragraph(title, title_style))
        
        # Subtitle
        if slide.get('subtitle'):
            elements.append(Paragraph(slide['subtitle'], subtitle_style))
        
        elements.append(Spacer(1, 20))
        
        # Bullet points
        bullet_points = slide.get('bullet_points', [])
        for point in bullet_points:
            bullet_text = f"• {point}"
            elements.append(Paragraph(bullet_text, bullet_style))
        
        # Speaker notes (optional)
        if include_notes and slide.get('speaker_notes'):
            elements.append(Spacer(1, 40))
            notes_style = ParagraphStyle('Notes', fontSize=10, textColor=rl_colors.gray, alignment=TA_LEFT)
            elements.append(Paragraph(f"📝 Notes: {slide['speaker_notes']}", notes_style))
    
    doc.build(elements)
    buffer.seek(0)
    return buffer.read()

@app.route('/generate-pdf', methods=['POST'])
def generate_pdf_presentation():
    """Generate PDF presentation from document content"""
    
    start_time = time.time()
    request_id = str(uuid.uuid4())[:8]
    
    try:
        data = request.json
        content = data.get('content', '')
        style = data.get('style', 'professional')
        slide_count = data.get('slide_count', data.get('slideCount', 10))
        title = data.get('title', 'AI Generated Presentation')

        structure_mode = data.get('structure', data.get('structureMode', 'classic'))
        structure_mode = (structure_mode or 'classic').strip().lower()
        topics_mode = structure_mode in ('topics', 'topic')
        include_notes = not topics_mode

        try:
            slide_count_int = int(slide_count)
        except Exception:
            slide_count_int = 10
        
        print(f'[{request_id}] PDF Request: slide_count={slide_count}, style={style}')
        
        if not content or len(content) < 100:
            return jsonify({'error': 'Content too short'}), 400
        
        # Generate slide structure
        if topics_mode:
            slides = generate_topic_slide_structure(content, slide_count_int, style)
            for s in slides:
                if isinstance(s, dict):
                    s.pop('speaker_notes', None)
        else:
            slides = generate_slide_structure(content, slide_count_int, style)
            # Ensure title slide has correct title
            if slides and slides[0].get('slide_type') == 'title':
                slides[0]['title'] = title
        
        # Generate PDF
        pdf_bytes = create_pdf_presentation(slides, style, include_notes=include_notes)
        
        # Save to temp file
        filename = f'presentation_{request_id}.pdf'
        filepath = os.path.join(tempfile.gettempdir(), filename)
        with open(filepath, 'wb') as f:
            f.write(pdf_bytes)
        
        duration = round(time.time() - start_time, 2)
        print(f'[{request_id}] PDF created in {duration}s')
        
        return jsonify({
            'success': True,
            'request_id': request_id,
            'download_url': f'/download-pdf/{request_id}',
            'filename': filename,
            'format': 'pdf',
            'slide_count': len(slides),
            'style': style,
            'structure': structure_mode,
            'include_notes': include_notes,
            'duration': duration
        })
        
    except Exception as e:
        print(f'[{request_id}] PDF Error: {e}')
        return jsonify({'error': str(e), 'request_id': request_id}), 500

@app.route('/generate-pdf-enhanced', methods=['POST'])
def generate_pdf_enhanced_presentation():
    """Generate enhanced PDF presentation with web search enrichment"""
    
    start_time = time.time()
    request_id = str(uuid.uuid4())[:8]
    
    try:
        data = request.json
        content = data.get('content', '')
        style = data.get('style', 'professional')
        slide_count = data.get('slide_count', data.get('slideCount', 10))
        title = data.get('title', 'AI Generated Presentation')

        structure_mode = data.get('structure', data.get('structureMode', 'classic'))
        structure_mode = (structure_mode or 'classic').strip().lower()
        topics_mode = structure_mode in ('topics', 'topic')
        include_notes = not topics_mode

        try:
            slide_count_int = int(slide_count)
        except Exception:
            slide_count_int = 10
        
        print(f'[{request_id}] PDF ENHANCED Request: slide_count={slide_count}, style={style}')
        
        if not content or len(content) < 100:
            return jsonify({'error': 'Content too short'}), 400
        
        all_web_results = []
        search_queries = []

        if topics_mode:
            print(f'[{request_id}] Topics mode: skipping web enrichment; using topic structure')
            slides = generate_topic_slide_structure(content, slide_count_int, style)
            for s in slides:
                if isinstance(s, dict):
                    s.pop('speaker_notes', None)
        else:
            # Step 1: Extract search queries
            print(f'[{request_id}] Step 1: Extracting search queries...')
            search_queries = extract_search_queries_from_content(content, slide_count_int)

            # Step 2: Perform web searches
            print(f'[{request_id}] Step 2: Searching web...')
            for query in search_queries[:5]:
                results = search_web_for_topic(query, max_results=3)
                all_web_results.extend(results)
                time.sleep(0.5)

            # Step 3: Generate enhanced structure
            print(f'[{request_id}] Step 3: Generating enhanced structure...')
            slides = generate_enhanced_slide_structure(content, slide_count_int, style, all_web_results)

            # Ensure professional opening and closing
            if slides and slides[0].get('slide_type') == 'title':
                slides[0]['title'] = title
                slides[0]['subtitle'] = 'Professional Presentation'

            if slides and slides[-1].get('slide_type') != 'closing':
                slides.append({
                    'slide_type': 'closing',
                    'title': 'Thank You',
                    'subtitle': 'Questions & Discussion',
                    'bullet_points': ['Thank you for your attention'],
                    'layout': 'no_image',
                    'speaker_notes': 'Thank the audience'
                })
        
        # Generate PDF
        pdf_bytes = create_pdf_presentation(slides, style, include_notes=include_notes)
        
        # Save to temp file
        filename = f'presentation_{request_id}.pdf'
        filepath = os.path.join(tempfile.gettempdir(), filename)
        with open(filepath, 'wb') as f:
            f.write(pdf_bytes)
        
        duration = round(time.time() - start_time, 2)
        print(f'[{request_id}] Enhanced PDF created in {duration}s')
        
        return jsonify({
            'success': True,
            'request_id': request_id,
            'download_url': f'/download-pdf/{request_id}',
            'filename': filename,
            'format': 'pdf',
            'slide_count': len(slides),
            'style': style,
            'structure': structure_mode,
            'include_notes': include_notes,
            'enhanced': (not topics_mode),
            'web_results_count': len(all_web_results),
            'search_queries': search_queries,
            'duration': duration
        })
        
    except Exception as e:
        print(f'[{request_id}] PDF Enhanced Error: {e}')
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e), 'request_id': request_id}), 500

@app.route('/download-pdf/<request_id>', methods=['GET'])
def download_pdf(request_id: str):
    """Download generated PDF presentation"""
    
    filename = f'presentation_{request_id}.pdf'
    filepath = os.path.join(tempfile.gettempdir(), filename)
    
    if not os.path.exists(filepath):
        return jsonify({'error': 'PDF not found or expired'}), 404
    
    return send_file(
        filepath,
        mimetype='application/pdf',
        as_attachment=True,
        download_name=filename
    )

# =============================================================================
# RUN SERVER
# =============================================================================

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, debug=False)
